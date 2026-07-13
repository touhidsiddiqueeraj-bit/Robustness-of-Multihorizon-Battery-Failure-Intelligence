from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY = RGBColor(0x0F, 0x17, 0x2A)
NAVY_LIGHT = RGBColor(0x1E, 0x29, 0x3B)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
CREAM = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
DANGER = RGBColor(0xDC, 0x26, 0x26)
SUCCESS = RGBColor(0x16, 0xA3, 0x4A)
WARN = RGBColor(0xD9, 0x77, 0x06)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gold_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(6858000-40000), Emu(12192000), Emu(40000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=18, color=TEXT, bold=False, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p

def add_label(slide, left, top, width, text):
    add_text_box(slide, left, top, width, Inches(0.3), text, font_size=11, color=GOLD, bold=True)

def add_slide_number(slide, num, total=11):
    add_text_box(slide, Inches(11.2), Inches(7.0), Inches(1.5), Inches(0.3), f"{num} / {total}", font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)

def add_figure(slide, path, left, top, width):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Emu(left), Emu(top), Emu(width))

def add_card(slide, left, top, width, height, num_text, label_text, sub_text="", border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.5)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(3)
    # number
    tf = add_text_box(slide, left + 40000, top + 20000, width - 80000, Inches(0.7), num_text, font_size=36, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    # label
    add_text_box(slide, left + 40000, top + height - 85000, width - 80000, Inches(0.35), label_text, font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    if sub_text:
        add_text_box(slide, left + 40000, top + height - 50000, width - 80000, Inches(0.25), sub_text, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

def add_zone(slide, left, top, width, height, title, ece, desc, bg_color, border_color, title_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    add_text_box(slide, left, top + Inches(0.3), width, Inches(0.5), title, font_size=36, color=title_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.9), width, Inches(0.4), ece, font_size=22, color=TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.3), width - Inches(0.4), Inches(0.6), desc, font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

fig_dir = os.path.join(os.path.dirname(__file__))

# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_gold_bar(slide)
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.8), "How Far Does the\nModel Hold?", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8), "Robustness of Multihorizon Battery\nFailure Intelligence Under Realistic Conditions", font_size=26, color=RGBColor(0xCC,0xCC,0xCC), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4), "Shikdar & Laaksonen — University of Vaasa", font_size=18, color=RGBColor(0x99,0x99,0x99), alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 1)

# ═══════════════════════════════════════════════
# SLIDE 2: Motivation
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_gold_bar(slide)
add_label(slide, Inches(0.8), Inches(0.4), Inches(4), "MOTIVATION")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6), "Lab results don't travel", font_size=36, color=NAVY, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), "The published hazard model was trained on clean lab data: full cycles, fixed C-rates, controlled temperatures.", font_size=18, color=TEXT)
add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(0.4), "Real BESS operation is fundamentally different:", font_size=18, color=TEXT)
items = [
    "Partial cycling — varying depths of discharge",
    "Temperature noise — sensor drift, thermal gradients",
    "Irregular rest — variable idle between cycles"
]
tf = add_text_box(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(2.5), "• " + items[0], font_size=18, color=TEXT)
for item in items[1:]:
    add_para(tf, "• " + item, font_size=18, color=TEXT, space_before=6)
# Warning box on right
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(Inches(7.5)), Emu(Inches(2.0)), Emu(Inches(4.5)), Emu(Inches(3.0)))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFE, 0xF2, 0xF2)
shape.line.color.rgb = RGBColor(0xFE, 0xCA, 0xCA)
shape.line.width = Pt(1)
add_text_box(slide, Inches(7.8), Inches(2.5), Inches(4), Inches(0.5), "Does calibration survive\nwhen the model meets real data?", font_size=22, color=DANGER, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 2)

# ═══════════════════════════════════════════════
# SLIDE 3: Approach (section)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_gold_bar(slide)
add_text_box(slide, Inches(1.5), Inches(0.5), Inches(10), Inches(0.3), "APPROACH", font_size=11, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(2.0), "Systematic robustness test\nacross four severity levels", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1.0), "Perturb lab data → apply frozen model →\nmeasure calibration → test recalibration", font_size=22, color=RGBColor(0xAA,0xAA,0xAA), alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 3)

# ═══════════════════════════════════════════════
# SLIDE 4: Perturbation methodology
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_gold_bar(slide)
add_label(slide, Inches(0.8), Inches(0.4), Inches(4), "EXPERIMENTAL DESIGN")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6), "Perturbation methodology", font_size=36, color=NAVY, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.2), Inches(0.5), "NASA discharge curves truncated at random DoD, with Gaussian temperature noise and proportional rest noise.", font_size=16, color=TEXT)
# Table
tbl_data = [
    ["Level", "DoD range", "Temp noise"],
    ["1", "75–100%", "±0.5°C"],
    ["2", "55–75%", "±1.0°C"],
    ["3", "35–55%", "±2.0°C"],
    ["4", "15–35%", "±3.0°C"],
]
rows, cols = 5, 3
table_shape = slide.shapes.add_table(rows, cols, Emu(Inches(0.8)), Emu(Inches(2.2)), Emu(Inches(5.2)), Emu(Inches(2.3)))
table = table_shape.table
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = tbl_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.name = "Calibri"
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = NAVY
            else:
                paragraph.font.color.rgb = TEXT
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.2), Inches(0.3), "5 seeds per level · 20 perturbed datasets · 20,560 records", font_size=13, color=TEXT_MUTED)
add_figure(slide, os.path.join(fig_dir, "fig_ece.png"), Inches(6.5), Inches(1.2), Inches(5.8))
add_slide_number(slide, 4)

# ═══════════════════════════════════════════════
# SLIDE 5: Calibration degrades
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_gold_bar(slide)
add_label(slide, Inches(0.8), Inches(0.4), Inches(4), "KEY FINDING")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), "Calibration degrades ~9×", font_size=36, color=NAVY, bold=True)
# Stat cards
cards = [
    ("0.031", "Clean baseline ECE", "H=20, lab conditions", SUCCESS),
    ("0.28", "Perturbed ECE", "Severity 1–4, H=20", DANGER),
    ("0.65–0.74", "Perturbed AUC", "Down from 0.985", WARN),
    ("S1 saturates", "Near-max loss", "Even 75–100% DoD", DANGER),
]
card_w = Inches(2.8)
card_h = Inches(1.6)
gap = Inches(0.2)
total_w = (Inches(0.8) * 2) + (card_w * 4) + (gap * 3)
start_x = Inches(0.8)
for i, (num, label, sub, color) in enumerate(cards):
    x = start_x + (card_w + gap) * i
    add_card(slide, int(x), int(Inches(1.5)), int(card_w), int(card_h), num, label, sub, color)

items = [
    "ECE jumps 0.031 → 0.277–0.293 — a nine-fold increase at H=20",
    "Degradation saturates at Severity 1: even minimal partial cycling causes near-max loss",
    "AUC drops more gracefully (0.985 → 0.65–0.74): discrimination ≠ calibration",
]
tf = add_text_box(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(3.0), "• " + items[0], font_size=18, color=TEXT)
for item in items[1:]:
    add_para(tf, "• " + item, font_size=18, color=TEXT, space_before=8)
add_slide_number(slide, 5)

# ═══════════════════════════════════════════════
# SLIDE 6: Root cause
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_gold_bar(slide)
add_label(slide, Inches(0.8), Inches(0.4), Inches(4), "DIAGNOSIS")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6), "Root cause: Vmin shift", font_size=36, color=NAVY, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5), "The primary driver is a shift in the minimum voltage feature under partial cycling.", font_size=18, color=TEXT)
add_card(slide, int(Inches(0.8)), int(Inches(2.2)), int(Inches(2.4)), int(Inches(1.4)), "2.3 V", "Full discharge", "", BORDER)
add_card(slide, int(Inches(3.5)), int(Inches(2.2)), int(Inches(2.4)), int(Inches(1.4)), "3.2 V", "75% DoD truncation", "", DANGER)
add_text_box(slide, Inches(0.8), Inches(3.9), Inches(5.5), Inches(0.8), "A 39% increase — the XGBoost model learned thresholds on clean ranges; the isotonic calibrator can't adapt to the shifted distribution.", font_size=16, color=TEXT_MUTED)
add_figure(slide, os.path.join(fig_dir, "fig_reliability.png"), Inches(6.5), Inches(1.2), Inches(5.8))
add_slide_number(slide, 6)

# ═══════════════════════════════════════════════
# SLIDE 7: SHAP evidence
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_gold_bar(slide)
add_label(slide, Inches(0.8), Inches(0.4), Inches(4), "EVIDENCE")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), "SHAP confirms Vmin dominance", font_size=36, color=NAVY, bold=True)
add_figure(slide, os.path.join(fig_dir, "fig_shap.png"), Inches(1.5), Inches(1.3), Inches(10.3))
add_text_box(slide, Inches(1.5), Inches(6.8), Inches(10), Inches(0.4), "Vmin is the dominant feature under both clean and perturbed regimes. Under perturbation, every feature distribution shifts, breaking the calibrator.", font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 7)

# ═══════════════════════════════════════════════
# SLIDE 8: Distribution panel
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_gold_bar(slide)
add_label(slide, Inches(0.8), Inches(0.4), Inches(4), "DIAGNOSIS")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), "Feature shift across severities", font_size=36, color=NAVY, bold=True)
add_figure(slide, os.path.join(fig_dir, "fig_dist.png"), Inches(1.0), Inches(1.4), Inches(11.3))
add_slide_number(slide, 9)

# ═══════════════════════════════════════════════
# SLIDE 10: Recalibration result (section)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_gold_bar(slide)
add_text_box(slide, Inches(1.5), Inches(0.5), Inches(10), Inches(0.3), "KEY RESULT", font_size=11, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(2.0), "Recalibration on 10%\nrecovers 71–88% of the gap", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8), "ECE drops from ~0.28 to 0.06–0.09\nEvaluated on held-out data", font_size=22, color=RGBColor(0xAA,0xAA,0xAA), alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 7)

# ═══════════════════════════════════════════════
# SLIDE 8: Recalibration detail
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_gold_bar(slide)
add_label(slide, Inches(0.8), Inches(0.4), Inches(4), "EVIDENCE")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6), "Recalibration: held-out evaluation", font_size=36, color=NAVY, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.2), Inches(0.4), "Isotonic regression fit on 10%, evaluated on held-out 90%:", font_size=16, color=TEXT)
tbl2_data = [
    ["Horizon", "Clean ECE", "Perturbed", "Recalibrated"],
    ["H=10", "0.010", "0.272–0.313", "0.057–0.076"],
    ["H=20", "0.031", "0.277–0.293", "0.060–0.073"],
    ["H=30", "0.013", "0.288–0.316", "0.076–0.092"],
    ["H=50", "0.023", "0.331–0.381", "0.080–0.094"],
]
rows2, cols2 = 5, 4
table2_shape = slide.shapes.add_table(rows2, cols2, Emu(Inches(0.8)), Emu(Inches(1.9)), Emu(Inches(5.2)), Emu(Inches(2.5)))
table2 = table2_shape.table
for r in range(rows2):
    for c in range(cols2):
        cell = table2.cell(r, c)
        cell.text = tbl2_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.name = "Calibri"
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = NAVY
            elif c == 3:
                paragraph.font.bold = True
                paragraph.font.color.rgb = SUCCESS
            else:
                paragraph.font.color.rgb = TEXT
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.2), Inches(0.3), "Recovery consistent across all severity levels and horizons", font_size=13, color=TEXT_MUTED)
add_figure(slide, os.path.join(fig_dir, "fig_combined.png"), Inches(6.5), Inches(1.2), Inches(5.8))
add_slide_number(slide, 10)

# ═══════════════════════════════════════════════
# SLIDE 11: Operating zones
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_gold_bar(slide)
add_label(slide, Inches(0.8), Inches(0.4), Inches(4), "FRAMEWORK")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), "Three operating zones", font_size=36, color=NAVY, bold=True)
zw = Inches(3.6)
zh = Inches(3.0)
zg = Inches(0.4)
zy = Inches(2.5)
add_zone(slide, int(Inches(0.8)), int(zy), int(zw), int(zh), "● Safe", "ECE < 0.05", "Clean lab conditions only", RGBColor(0xF0,0xFD,0xF4), RGBColor(0xBB,0xF7,0xD0), SUCCESS)
add_zone(slide, int(Inches(0.8+3.6+0.4)), int(zy), int(zw), int(zh), "● Warning", "ECE 0.05–0.10", "Recalibrated operation\n10% operational sample", RGBColor(0xFF,0xFB,0xEB), RGBColor(0xFD,0xE6,0x8A), WARN)
add_zone(slide, int(Inches(0.8+7.2+0.8)), int(zy), int(zw), int(zh), "● Unsafe", "ECE > 0.10", "Direct deployment under\nany partial cycling (ECE > 0.27)", RGBColor(0xFE,0xF2,0xF2), RGBColor(0xFE,0xCA,0xCA), DANGER)
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4), "Direct deployment under partial cycling pushes ECE > 0.27 — firmly in the unsafe zone", font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 11)

# ═══════════════════════════════════════════════
# SLIDE 12: Deployment path
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_gold_bar(slide)
add_label(slide, Inches(0.8), Inches(0.4), Inches(4), "ACTION")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), "Practical deployment path", font_size=36, color=NAVY, bold=True)
steps = [
    "1. Deploy the frozen XGBoost model as-is",
    "2. Collect 50–100 operational cycles from normal BESS operation",
    "3. Fit a new isotonic regression on these cycles",
    "4. Deploy the recalibrated model — no base classifier retraining needed",
]
tf = add_text_box(slide, Inches(0.8), Inches(1.6), Inches(7.5), Inches(3.0), steps[0], font_size=20, color=TEXT)
for step in steps[1:]:
    add_para(tf, step, font_size=20, color=TEXT, space_before=10)
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(7.5), Inches(0.4), "No original training data required · No additional feature engineering", font_size=14, color=TEXT_MUTED)
# Highlight box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(Inches(9.0)), Emu(Inches(2.0)), Emu(Inches(3.5)), Emu(Inches(2.0)))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = BORDER
shape.line.width = Pt(1)
add_text_box(slide, Inches(9.2), Inches(2.2), Inches(3.1), Inches(0.6), "10%", font_size=48, color=SUCCESS, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.2), Inches(2.8), Inches(3.1), Inches(0.4), "operational sample suffices", font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 12)

# ═══════════════════════════════════════════════
# SLIDE 13: Conclusion
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_gold_bar(slide)
add_text_box(slide, Inches(1.5), Inches(0.5), Inches(10), Inches(0.3), "TAKEAWAY", font_size=11, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(2.0), "Calibration is fragile.\nRecalibration is practical.", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1.2), "Multihorizon hazard models are not inherently robust to\noperational distribution shift, but a lightweight recalibration\nstep suffices for field deployment.", font_size=22, color=RGBColor(0xAA,0xAA,0xAA), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4), "Shikdar & Laaksonen — University of Vaasa, Finland", font_size=16, color=RGBColor(0x88,0x88,0x88), alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 11)

out = os.path.join(fig_dir, "robustness_presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
